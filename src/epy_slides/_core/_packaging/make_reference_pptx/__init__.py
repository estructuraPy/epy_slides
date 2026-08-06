"""Generate per-theme PowerPoint reference decks from the epyson themes.

Pandoc's ``pptx`` writer copies the master slide layouts, fonts and colours
from a ``--reference-doc``. python-pptx's default template already carries
the seven standard layouts Pandoc matches by name (Title Slide, Section
Header, Two Content, Comparison, Content with Caption, Title and Content,
Blank); this script only re-skins each one's *theme part* — the colour
scheme and the major/minor fonts — from a :class:`Theme`'s ``css_vars`` so
the exported deck carries the same identity as the live preview.

Run from the project root::

    python src/epy_slides/_core/_packaging/make_reference_pptx/__init__.py

It writes ``src/epy_slides/_config/_assets/reference_pptx/<theme>.pptx`` for
every bundled theme (the build-time dependency ``python-pptx`` is not
required at runtime — Pandoc consumes the generated files).
"""

from __future__ import annotations

import sys
from pathlib import Path

# Repo root: this module lives at src/epy_slides/_core/_packaging/
# make_reference_pptx/__init__.py, so it takes 5 parents to reach the
# root (make_reference_pptx -> _packaging -> _core -> epy_slides -> src
# -> root). parents[4] dates from when this was a flat .py and silently
# wrote the decks into a ghost src/src/ tree.
ROOT = Path(__file__).resolve().parents[5]
sys.path.insert(0, str(ROOT / "src"))

from lxml import etree  # noqa: E402
from pptx import Presentation  # noqa: E402
from pptx.opc.constants import RELATIONSHIP_TYPE as RT  # noqa: E402
from pptx.oxml.ns import qn  # noqa: E402
from pptx.util import Emu  # noqa: E402

# PowerPoint 16:9 "widescreen" (13.333 x 7.5 in). python-pptx's default
# template is 4:3 (10 x 7.5 in); the decks are 16:9, so a 4:3 reference makes
# Pandoc lay 16:9 content into a 4:3 frame and the text shifts / clips.
# A `<theme>_43.pptx` variant keeps the native 4:3 canvas for decks whose
# front matter declares `aspect-ratio: 4:3`.
_WIDESCREEN_W = Emu(12192000)
_WIDESCREEN_H = Emu(6858000)

# Master text-style sizes (pt*100). python-pptx's template ships Office's
# print-era defaults (body 32/28/24 pt), which fill the frame after a few
# bullets and overflow; these track the live reveal.js preview instead
# (~21 pt body on the 13.33 in canvas), so the exported deck reads like
# the on-screen one and holds comparable content per slide.
_TITLE_SIZE = 3600
_BODY_SIZES = [2200, 2000, 1800, 1600, 1600, 1600, 1600, 1600, 1600]

from epy_slides._core import themes  # noqa: E402
from epy_slides._core.themes_base import Theme  # noqa: E402

OUT_DIR = ROOT / "src" / "epy_slides" / "_config" / "_assets" / "reference_pptx"

# epyson css var -> theme colour-scheme slot (PowerPoint names).
_COLOR_MAP = {
    "dk1": "fg",
    "lt1": "bg",
    "dk2": "heading-color",
    "lt2": "bg-soft",
    "accent1": "link",
    "accent2": "link-hover",
    "accent3": "mark-bg",
    "accent4": "table-header-bg",
    "accent5": "border",
    "accent6": "quote-rule",
    "hlink": "link",
    "folHlink": "link-hover",
}


def _hex(value: str, default: str = "000000") -> str:
    """Return a 6-digit RRGGBB hex (no ``#``) from a css colour value."""
    value = (value or "").strip().lstrip("#")
    return value.upper() if len(value) == 6 else default


def _first_font(value: str, default: str) -> str:
    """Return the first family name from a CSS font-family list."""
    if not value:
        return default
    first = value.split(",")[0].strip().strip('"').strip("'")
    return first or default


def _set_scheme_color(clr_scheme, slot: str, hex_value: str) -> None:
    """Set a colour-scheme slot to a solid ``srgbClr``."""
    el = clr_scheme.find(qn(f"a:{slot}"))
    if el is None:
        return
    for child in list(el):
        el.remove(child)
    srgb = etree.SubElement(el, qn("a:srgbClr"))
    srgb.set("val", hex_value)


def _set_font(font_scheme, which: str, typeface: str) -> None:
    """Set the latin typeface of ``majorFont`` / ``minorFont``."""
    family = font_scheme.find(qn(f"a:{which}"))
    if family is None:
        return
    latin = family.find(qn("a:latin"))
    if latin is not None:
        latin.set("typeface", typeface)


def _widen_placeholders(container, ratio: float) -> None:
    """Scale the x position and width of placeholders that own their geometry.

    Converts a 4:3 layout/master to fill the wider 16:9 canvas: only the
    horizontal axis grows (16:9 height == 4:3 height), so the ``a:off/@x`` and
    ``a:ext/@cx`` of each *explicit* ``a:xfrm`` scale; top/height stay put.
    Placeholders that INHERIT their geometry have no ``a:xfrm`` and are skipped
    — they would otherwise be scaled twice (once on the master they inherit
    from, once here on the resolved value), blowing a 9 in box up to 16 in.
    """
    for ph in container.placeholders:
        sppr = ph._element.spPr
        xfrm = sppr.find(qn("a:xfrm")) if sppr is not None else None
        if xfrm is None:
            continue
        off = xfrm.find(qn("a:off"))
        ext = xfrm.find(qn("a:ext"))
        if off is not None and off.get("x") is not None:
            off.set("x", str(int(int(off.get("x")) * ratio)))
        if ext is not None and ext.get("cx") is not None:
            ext.set("cx", str(int(int(ext.get("cx")) * ratio)))


def _retune_text_styles(master) -> None:
    """Replace the master's title/body default sizes with preview-like ones."""
    tx_styles = master._element.find(qn("p:txStyles"))
    if tx_styles is None:
        return
    title_style = tx_styles.find(qn("p:titleStyle"))
    if title_style is not None:
        for pr in title_style.iter(qn("a:defRPr")):
            if pr.get("sz") is not None:
                pr.set("sz", str(_TITLE_SIZE))
    body_style = tx_styles.find(qn("p:bodyStyle"))
    if body_style is not None:
        for lvl, size in enumerate(_BODY_SIZES, start=1):
            pr_parent = body_style.find(qn(f"a:lvl{lvl}pPr"))
            if pr_parent is None:
                continue
            pr = pr_parent.find(qn("a:defRPr"))
            if pr is not None and pr.get("sz") is not None:
                pr.set("sz", str(size))


def _enable_autofit(container) -> None:
    """Set "shrink text on overflow" on content/title placeholder bodies.

    Pandoc rewrites each slide shape's own ``bodyPr``, so this does not
    reach generated slides (``_pptx_polish`` handles those at export
    time) — but decks edited by hand from these masters keep their text
    inside the frame.
    """
    chrome = {"dt", "ftr", "sldNum"}
    for ph in container.placeholders:
        ph_el = ph._element.find(
            f"{qn('p:nvSpPr')}/{qn('p:nvPr')}/{qn('p:ph')}"
        )
        if ph_el is not None and (ph_el.get("type") or "body") in chrome:
            continue
        body_pr = ph.text_frame._txBody.bodyPr
        for tag in ("a:normAutofit", "a:spAutoFit", "a:noAutofit"):
            for old in body_pr.findall(qn(tag)):
                body_pr.remove(old)
        etree.SubElement(body_pr, qn("a:normAutofit"))


def _drop_stale_size_type(prs) -> None:
    """Remove ``p:sldSz/@type`` after resizing (stale "screen4x3" label)."""
    sld_sz = prs._element.find(qn("p:sldSz"))
    if sld_sz is not None and sld_sz.get("type") is not None:
        del sld_sz.attrib["type"]


def build_reference(
    theme: Theme, target: Path, *, widescreen: bool = True
) -> None:
    """Write a themed reference deck for ``theme`` to ``target``.

    ``widescreen=True`` produces the 16:9 deck (13.333 x 7.5 in);
    ``False`` keeps the template's native 4:3 canvas (10 x 7.5 in).
    """
    prs = Presentation()
    if widescreen:
        # Switch the default 4:3 canvas to 16:9 widescreen and widen the
        # layout placeholders to fill it, so Pandoc lays slide content
        # across the full 16:9 width instead of a 4:3 column.
        ratio = float(_WIDESCREEN_W) / float(prs.slide_width)
        prs.slide_width = _WIDESCREEN_W
        prs.slide_height = _WIDESCREEN_H
        for master in prs.slide_masters:
            _widen_placeholders(master, ratio)
            for layout in master.slide_layouts:
                _widen_placeholders(layout, ratio)
        _drop_stale_size_type(prs)
    for master in prs.slide_masters:
        _retune_text_styles(master)
        _enable_autofit(master)
        for layout in master.slide_layouts:
            _enable_autofit(layout)
    css = theme.css_vars
    for master in prs.slide_masters:
        theme_part = master.part.part_related_by(RT.THEME)
        theme_el = etree.fromstring(theme_part.blob)
        elements = theme_el.find(qn("a:themeElements"))
        clr_scheme = elements.find(qn("a:clrScheme"))
        font_scheme = elements.find(qn("a:fontScheme"))
        for slot, var in _COLOR_MAP.items():
            _set_scheme_color(clr_scheme, slot, _hex(css.get(var, "")))
        _set_font(
            font_scheme, "majorFont",
            _first_font(css.get("font-family-headings", ""), "Calibri Light"),
        )
        _set_font(
            font_scheme, "minorFont",
            _first_font(css.get("font-family-text", ""), "Calibri"),
        )
        theme_part._blob = etree.tostring(
            theme_el, xml_declaration=True, encoding="UTF-8", standalone=True
        )
    prs.save(str(target))


def main() -> int:
    """Generate a reference deck for every bundled theme."""
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    init = OUT_DIR / "__init__.py"
    if not init.exists():
        init.write_text(
            '"""Per-theme PowerPoint reference decks (Pandoc pptx)."""\n',
            encoding="utf-8",
        )
    for theme_id, theme in themes.THEMES.items():
        target = OUT_DIR / f"{theme_id}.pptx"
        build_reference(theme, target, widescreen=True)
        print(f"wrote {target.relative_to(ROOT)}")
        target_43 = OUT_DIR / f"{theme_id}_43.pptx"
        build_reference(theme, target_43, widescreen=False)
        print(f"wrote {target_43.relative_to(ROOT)}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
